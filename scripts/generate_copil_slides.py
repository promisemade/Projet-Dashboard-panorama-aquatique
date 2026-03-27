from __future__ import annotations

import json
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DASHBOARD_FILE = ROOT / "public" / "data" / "dashboard.json"
OUTPUT_FILE = ROOT / f"copil_filiere_aquatique_{date.today().isoformat()}.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(0x00, 0x00, 0x91)
BLUE_LIGHT = RGBColor(0xE8, 0xEE, 0xFF)
RED = RGBColor(0xC9, 0x19, 0x1E)
RED_LIGHT = RGBColor(0xFE, 0xF4, 0xF4)
GREEN = RGBColor(0x18, 0x75, 0x3C)
GREEN_LIGHT = RGBColor(0xE6, 0xF7, 0xED)
ORANGE = RGBColor(0xB3, 0x40, 0x00)
ORANGE_LIGHT = RGBColor(0xFE, 0xF0, 0xE8)
GOLD = RGBColor(0x7A, 0x58, 0x00)
GOLD_LIGHT = RGBColor(0xFE, 0xF7, 0xDA)
GRAY_975 = RGBColor(0x16, 0x16, 0x16)
GRAY_425 = RGBColor(0x66, 0x66, 0x66)
GRAY_200 = RGBColor(0xDD, 0xDD, 0xDD)
GRAY_100 = RGBColor(0xF6, 0xF6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def fmt_int(value: int | float) -> str:
    return f"{int(round(value)):,}".replace(",", " ")


def fmt_percent(value: float, digits: int = 1) -> str:
    return f"{value * 100:.{digits}f}".replace(".", ",") + " %"


def rgb_from(color: RGBColor) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def load_metrics() -> dict[str, str]:
    payload = json.loads(DASHBOARD_FILE.read_text(encoding="utf-8"))
    overview = payload["overview"]
    school = payload["school_demand_overview"]
    accessibility = payload["accessibility_overview"]
    transit = payload["transit_overview"]
    statuses = Counter(row["operational_status_code"] for row in payload["installation_status"])
    projects_count = len(payload["projects_in_progress"])

    closed_or_works = statuses.get("temporary_closed", 0) + statuses.get("closed", 0)
    seasonal_or_verify = statuses.get("seasonal", 0) + statuses.get("verify", 0)

    return {
        "population_total": fmt_int(overview["population_total"]),
        "installations_total": fmt_int(overview["installations_total"]),
        "bassins_total": fmt_int(overview["bassins_total"]),
        "surface_total": fmt_int(overview["surface_totale_bassins_m2"]) + " m²",
        "licences_2024": fmt_int(overview["licences_ffn_2024"]),
        "communes_without_basin": fmt_int(overview["communes_avec_licences_sans_bassin"]),
        "communes_without_basin_share": fmt_percent(
            overview["communes_avec_licences_sans_bassin"] / overview["communes_total"]
        ),
        "regie_share": fmt_percent(overview["bassins_regie"] / overview["bassins_total"]),
        "school_sites": fmt_int(school["schools_total"]),
        "students_total": fmt_int(school["students_total"]),
        "school_bassins": fmt_int(overview["bassins_usage_scolaires"]),
        "school_bassins_share": fmt_percent(overview["bassins_usage_scolaires"] / overview["bassins_total"]),
        "students_per_installation": fmt_int(school["students_per_installation"]),
        "students_per_school_basin": fmt_int(school["students_per_school_basin"]),
        "students_within_15min": fmt_percent(school["students_within_15min_installation_share"]),
        "population_within_15min": fmt_percent(accessibility["population_within_15min_share"]),
        "population_within_500m_tc": fmt_percent(transit["population_within_500m_share"]),
        "closed_or_works": fmt_int(closed_or_works),
        "seasonal_or_verify": fmt_int(seasonal_or_verify),
        "projects_count": fmt_int(projects_count),
    }


def set_text_frame_style(text_frame, font_name: str = "Arial", size: int = 18, color: RGBColor = GRAY_975) -> None:
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.color.rgb = rgb_from(color)


def add_header(slide, kicker: str, title: str, subtitle: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.58))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()

    kicker_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(5.0), Inches(0.25))
    tf = kicker_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = kicker
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.88), Inches(8.6), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = GRAY_975

    subtitle_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.42), Inches(11.8), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY_425


def add_footer(slide) -> None:
    footer_box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.25))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Sources : Data.Sports, Data.Education, geo.api.gouv.fr, OpenStreetMap/OSRM, transport.data.gouv.fr"
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY_425


def add_stat_card(slide, x: float, y: float, w: float, h: float, label: str, value: str, detail: str, accent: RGBColor, bg: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = accent
    card.line.width = Pt(1.2)

    label_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.22))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = accent

    value_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.42), w - Inches(0.36), Inches(0.42))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = GRAY_975

    detail_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.92), w - Inches(0.36), h - Inches(1.0))
    tf = detail_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = detail
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_425


def add_message_box(slide, title: str, body: str, x: float = 0.55, y: float = 5.95, w: float = 12.2, h: float = 0.75) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_100
    box.line.color.rgb = GRAY_200
    box.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(2.0), Inches(0.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = BLUE

    body_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.32), Inches(w - 0.4), Inches(h - 0.34))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Arial"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_975


def build_slide_1(prs: Presentation, m: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Comité de pilotage · Enquête besoins de la filière aquatique",
        "1. Un parc régional réel, mais une couverture territoriale inégale",
        "Les données publiques croisées montrent une offre existante, mais très inégalement répartie sur le territoire.",
    )
    cards = [
        ("Installations", m["installations_total"], "Sites aquatiques recensés dans le périmètre régional.", BLUE, BLUE_LIGHT),
        ("Bassins", m["bassins_total"], "Bassins identifiés dans le socle régional exploité.", ORANGE, ORANGE_LIGHT),
        ("Surface totale", m["surface_total"], "Surface cumulée des bassins renseignés.", GREEN, GREEN_LIGHT),
        ("Licences FFN 2024", m["licences_2024"], "Pratique fédérée repérée dans les données publiques.", BLUE, BLUE_LIGHT),
        (
            "Communes sans bassin",
            m["communes_without_basin"],
            f"{m['communes_without_basin_share']} des communes ont des licences mais aucun bassin.",
            RED,
            RED_LIGHT,
        ),
        ("Régie publique", m["regie_share"], "Part des bassins gérés en régie publique.", GOLD, GOLD_LIGHT),
    ]
    positions = [
        (0.55, 2.0),
        (4.5, 2.0),
        (8.45, 2.0),
        (0.55, 3.85),
        (4.5, 3.85),
        (8.45, 3.85),
    ]
    for (x, y), card in zip(positions, cards, strict=True):
        add_stat_card(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.5), *card)
    add_message_box(
        slide,
        "À retenir",
        "Le sujet n'est pas seulement le volume d'offre. Le principal enjeu est la continuité territoriale du service aquatique dans une région de près de 6 millions d'habitants.",
    )
    add_footer(slide)


def build_slide_2(prs: Presentation, m: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Comité de pilotage · Enquête besoins de la filière aquatique",
        "2. Le scolaire structure une large part du besoin",
        "Le croisement entre équipements et données scolaires met en évidence une pression potentielle forte sur les installations.",
    )
    cards = [
        ("Établissements", m["school_sites"], "Écoles, collèges et lycées intégrés à l'analyse.", BLUE, BLUE_LIGHT),
        ("Élèves", m["students_total"], "Effectifs scolaires repérés dans le périmètre régional.", ORANGE, ORANGE_LIGHT),
        ("Bassins scolaires", m["school_bassins"], f"Soit {m['school_bassins_share']} du parc recensé.", GREEN, GREEN_LIGHT),
        ("Élèves / installation", m["students_per_installation"], "Pression scolaire potentielle rapportée au nombre de sites.", RED, RED_LIGHT),
        ("Élèves / bassin scolaire", m["students_per_school_basin"], "Lecture plus fine sur les bassins à usage scolaire.", GOLD, GOLD_LIGHT),
        ("Élèves à < 15 min", m["students_within_15min"], "Accès voiture estimé vers l'installation la plus proche.", BLUE, BLUE_LIGHT),
    ]
    positions = [
        (0.55, 2.0),
        (4.5, 2.0),
        (8.45, 2.0),
        (0.55, 3.85),
        (4.5, 3.85),
        (8.45, 3.85),
    ]
    for (x, y), card in zip(positions, cards, strict=True):
        add_stat_card(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.5), *card)
    add_message_box(
        slide,
        "À retenir",
        "Le scolaire n'est pas un usage marginal. Il faut raisonner les besoins non seulement en nombre d'équipements, mais en capacité réelle d'accueil scolaire.",
    )
    add_footer(slide)


def build_slide_3(prs: Presentation, m: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Comité de pilotage · Enquête besoins de la filière aquatique",
        "3. Accessibilité globalement correcte, mais parc fragilisé",
        "L'accès en voiture reste globalement bon, mais le parc est marqué par des fermetures, des travaux et des besoins de renouvellement.",
    )
    cards = [
        ("Population à < 15 min", m["population_within_15min"], "Part de la population à moins de 15 min en voiture d'une installation.", BLUE, BLUE_LIGHT),
        ("Population à < 500 m TC", m["population_within_500m_tc"], "Proximité d'un arrêt ou d'une gare, lecture plus exigeante.", ORANGE, ORANGE_LIGHT),
        ("Fermées / travaux", m["closed_or_works"], "Installations fermées, hors service ou temporairement indisponibles.", RED, RED_LIGHT),
        ("Saisonnières / à vérifier", m["seasonal_or_verify"], "Cas encore fragiles ou à confirmer dans la couche d'exploitation.", GOLD, GOLD_LIGHT),
        ("Projets en cours", m["projects_count"], "Constructions, réhabilitations lourdes ou projets encore incertains.", GREEN, GREEN_LIGHT),
    ]
    positions = [
        (0.55, 2.0),
        (4.5, 2.0),
        (8.45, 2.0),
        (2.52, 3.85),
        (6.47, 3.85),
    ]
    sizes = [
        (3.55, 1.5),
        (3.55, 1.5),
        (3.55, 1.5),
        (3.55, 1.5),
        (3.55, 1.5),
    ]
    for (x, y), (w, h), card in zip(positions, sizes, cards, strict=True):
        add_stat_card(slide, Inches(x), Inches(y), Inches(w), Inches(h), *card)
    add_message_box(
        slide,
        "À retenir",
        "Les projets repérés sont importants, mais ils ne suffiront pas à eux seuls. La question porte aussi sur la sécurisation, la réouverture et la remise à niveau de l'offre existante.",
    )
    add_footer(slide)


def main() -> None:
    metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    build_slide_1(prs, metrics)
    build_slide_2(prs, metrics)
    build_slide_3(prs, metrics)

    prs.save(OUTPUT_FILE)
    print(f"PowerPoint written to {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
